"""
Office document parser for the Vectara vector database.

Supports Microsoft Office formats (Word, Excel, PowerPoint).
"""

from typing import Any, Dict, List, Union
from pathlib import Path

from .base import DocumentParser
from ..exceptions import ParsingError


class OfficeParser(DocumentParser):
    """
    Parser for Microsoft Office documents.
    
    Supports .docx, .xlsx, .pptx files using python-docx, openpyxl, and python-pptx.
    """
    
    def __init__(self, config: Dict[str, Any]):
        """Initialize Office parser."""
        super().__init__(config)
        self.extract_tables = config.get("extract_tables", True)
        self.extract_images = config.get("extract_images", False)
        self.chunk_size = config.get("chunk_size", 1000)
        self.chunk_overlap = config.get("chunk_overlap", 100)
        self.max_file_size = config.get("max_file_size", 100 * 1024 * 1024)  # 100MB
    
    def initialize(self) -> None:
        """Initialize the Office parser."""
        try:
            # Try to import required libraries
            import docx
            import openpyxl
            import pptx
            self._initialized = True
        except ImportError as e:
            raise ParsingError(f"Required libraries not installed: {str(e)}")
    
    def parse(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Parse an Office document."""
        if not self._initialized:
            raise ParsingError("Parser not initialized")
        
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise ParsingError(f"File not found: {file_path}")
        
        if not self.can_parse(file_path):
            raise ParsingError(f"Unsupported file format: {file_path.suffix}")
        
        try:
            # Check file size
            file_size = file_path.stat().st_size
            if file_size > self.max_file_size:
                raise ParsingError(f"File too large: {file_size} bytes")
            
            suffix = file_path.suffix.lower()
            
            if suffix == ".docx":
                return self._parse_docx(file_path, file_size)
            elif suffix in [".xlsx", ".xls"]:
                return self._parse_excel(file_path, file_size)
            elif suffix == ".pptx":
                return self._parse_pptx(file_path, file_size)
            else:
                raise ParsingError(f"Unsupported Office format: {suffix}")
                
        except Exception as e:
            raise ParsingError(f"Failed to parse Office document {file_path}: {str(e)}")
    
    def _parse_docx(self, file_path: Path, file_size: int) -> Dict[str, Any]:
        """Parse Word document."""
        import docx
        
        doc = docx.Document(file_path)
        
        # Extract text from paragraphs
        text_parts = []
        paragraphs_info = []
        
        for i, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if text:
                text_parts.append(text)
                paragraphs_info.append({
                    "paragraph_number": i + 1,
                    "text_length": len(text),
                    "style": paragraph.style.name if paragraph.style else "Normal"
                })
        
        full_text = "\n".join(text_parts)
        
        # Extract tables if requested
        tables_info = []
        if self.extract_tables:
            for table_num, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables_info.append({
                    "table_number": table_num + 1,
                    "rows": len(table_data),
                    "columns": len(table_data[0]) if table_data else 0,
                    "data": table_data
                })
        
        metadata = {
            "file_name": file_path.name,
            "file_size": file_size,
            "file_type": "docx",
            "num_paragraphs": len(paragraphs_info),
            "num_tables": len(tables_info),
            "paragraphs_info": paragraphs_info,
            "tables_info": tables_info
        }
        
        # Create chunks
        chunks = self._create_chunks(full_text, paragraphs_info)
        
        structure = {
            "type": "word_document",
            "paragraphs": len(paragraphs_info),
            "tables": len(tables_info),
            "estimated_reading_time": len(full_text.split()) // 200
        }
        
        return {
            "text": full_text,
            "metadata": metadata,
            "chunks": chunks,
            "structure": structure
        }
    
    def _parse_excel(self, file_path: Path, file_size: int) -> Dict[str, Any]:
        """Parse Excel document."""
        import openpyxl
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        text_parts = []
        sheets_info = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                
                if row_text:
                    sheet_text.append(" | ".join(row_text))
            
            if sheet_text:
                sheet_content = "\n".join(sheet_text)
                text_parts.append(f"Sheet: {sheet_name}\n{sheet_content}")
                sheets_info.append({
                    "sheet_name": sheet_name,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column,
                    "text_length": len(sheet_content)
                })
        
        full_text = "\n\n".join(text_parts)
        
        metadata = {
            "file_name": file_path.name,
            "file_size": file_size,
            "file_type": "excel",
            "num_sheets": len(workbook.sheetnames),
            "sheet_names": workbook.sheetnames,
            "sheets_info": sheets_info
        }
        
        chunks = self._create_chunks(full_text, sheets_info)
        
        structure = {
            "type": "excel_spreadsheet",
            "sheets": len(workbook.sheetnames),
            "total_cells": sum(s["max_row"] * s["max_column"] for s in sheets_info)
        }
        
        return {
            "text": full_text,
            "metadata": metadata,
            "chunks": chunks,
            "structure": structure
        }
    
    def _parse_pptx(self, file_path: Path, file_size: int) -> Dict[str, Any]:
        """Parse PowerPoint presentation."""
        import pptx
        
        presentation = pptx.Presentation(file_path)
        
        text_parts = []
        slides_info = []
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slide_content = "\n".join(slide_text)
                text_parts.append(f"Slide {slide_num + 1}: {slide_content}")
                slides_info.append({
                    "slide_number": slide_num + 1,
                    "text_length": len(slide_content),
                    "shapes_with_text": len(slide_text)
                })
        
        full_text = "\n\n".join(text_parts)
        
        metadata = {
            "file_name": file_path.name,
            "file_size": file_size,
            "file_type": "pptx",
            "num_slides": len(presentation.slides),
            "slides_info": slides_info
        }
        
        chunks = self._create_chunks(full_text, slides_info)
        
        structure = {
            "type": "powerpoint_presentation",
            "slides": len(presentation.slides),
            "estimated_reading_time": len(full_text.split()) // 200
        }
        
        return {
            "text": full_text,
            "metadata": metadata,
            "chunks": chunks,
            "structure": structure
        }
    
    def get_supported_formats(self) -> List[str]:
        """Get supported Office file formats."""
        return [".docx", ".xlsx", ".xls", ".pptx"]
    
    def can_parse(self, file_path: Union[str, Path]) -> bool:
        """Check if file can be parsed."""
        file_path = Path(file_path)
        return file_path.suffix.lower() in self.get_supported_formats()
    
    def _create_chunks(self, text: str, structure_info: List[Dict]) -> List[Dict[str, Any]]:
        """Create text chunks with structure information."""
        chunks = []
        start = 0
        chunk_id = 0
        
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            
            # Try to break at logical boundaries
            if end < len(text):
                last_newline = text.rfind('\n', start, end)
                if last_newline > start + self.chunk_size // 2:
                    end = last_newline + 1
            
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "id": chunk_id,
                    "text": chunk_text,
                    "start_pos": start,
                    "end_pos": end,
                    "metadata": {
                        "chunk_type": "office_segment"
                    }
                })
                chunk_id += 1
            
            start = max(end - self.chunk_overlap, start + 1)
        
        return chunks
